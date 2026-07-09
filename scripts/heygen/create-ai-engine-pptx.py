#!/usr/bin/env python3
"""
Gas Station COO — AI Engine Implementation Deck (companion to developer architecture PPT)

Source: gas_station_ai_engine_prd.html, gas_station_site_feasibility.md,
        gas_station_video_scripts_developer.html (engine episodes)

Focus: components, development, deployment, training, forecasting, hosting,
       and how the engine handles multi-store / multi-tenant volume.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# palette — purple/teal accent for engine deck (distinct from platform blue deck)
NAVY = RGBColor(0x0D, 0x3B, 0x66)
BLUE = RGBColor(0x00, 0x66, 0xCC)
BLUE_LT = RGBColor(0xE8, 0xF4, 0xF8)
TEAL = RGBColor(0x0D, 0x94, 0x88)
TEAL_LT = RGBColor(0xE0, 0xF7, 0xF4)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
ORANGE_LT = RGBColor(0xFE, 0xF5, 0xE7)
PURPLE = RGBColor(0x6C, 0x34, 0x8D)
PURPLE_LT = RGBColor(0xF5, 0xEE, 0xF8)
GREEN = RGBColor(0x1E, 0x84, 0x4E)
GREEN_LT = RGBColor(0xE8, 0xF8, 0xEF)
GRAY = RGBColor(0x4A, 0x55, 0x68)
GRAY_LT = RGBColor(0xF4, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x23, 0x32)
ACCENT_LINE = RGBColor(0x94, 0xA3, 0xB8)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUT = Path(__file__).resolve().parents[2] / (
    "documentation/tenant_management/gas_station_site/videos/generated/decks/"
    "gas-station-coo-ai-engine-implementation.pptx"
)

# Import diagram generator (same folder)
import sys

sys.path.insert(0, str(Path(__file__).resolve().parent))
from ai_engine_excalidraw_diagrams import (  # noqa: E402
    CONNECTIVITY_DIAGRAM_SLIDE_META,
    DIAGRAM_DIR,
    DIAGRAM_SLIDE_META,
    generate_all,
)


def set_run(run, size=18, bold=False, color=DARK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_text_box(slide, left, top, width, height, text, *, size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box


def add_rounded_rect(slide, left, top, width, height, fill, *, line=None, line_w=1.25):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    try:
        shape.adjustments[0] = 0.15
    except Exception:
        pass
    return shape


def label_shape(shape, text, *, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    try:
        tf._txBody.get_or_add_bodyPr().set(qn("a:anchor"), "ctr" if valign == MSO_ANCHOR.MIDDLE else "t")
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)


def add_bullet_box(slide, left, top, width, height, title, bullets, *, fill=TEAL_LT, title_color=TEAL, body_color=DARK):
    card = add_rounded_rect(slide, left, top, width, height, fill, line=title_color, line_w=1)
    tf = card.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_run(run, size=15, bold=True, color=title_color)
    for b in bullets:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"• {b}"
        set_run(run, size=13, color=body_color)
    return card


def add_arrow(slide, x1, y1, x2, y2, color=TEAL):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(2.25)


def slide_header(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PURPLE
    bar.line.fill.background()
    add_text_box(slide, Inches(0.4), Inches(0.18), Inches(12.2), Inches(0.45), title, size=26, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, Inches(0.4), Inches(0.55), Inches(12.2), Inches(0.35), subtitle, size=13, color=PURPLE_LT)
    add_text_box(
        slide,
        Inches(0.4),
        Inches(7.15),
        Inches(12.4),
        Inches(0.3),
        "Gas Station AI Engine · Implementation deck · gas_station_ai_engine_prd.html",
        size=10,
        color=ACCENT_LINE,
    )


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_diagram_slide(prs, title: str, subtitle: str, png_path: Path):
    """One slide with header + embedded diagram PNG (fits 16:9 content area)."""
    s = blank_slide(prs)
    slide_header(s, title, subtitle)
    if png_path.is_file():
        s.shapes.add_picture(str(png_path), Inches(0.35), Inches(1.05), width=Inches(12.6))
    else:
        add_text_box(
            s,
            Inches(0.5),
            Inches(2.5),
            Inches(12),
            Inches(1),
            f"Diagram not found: {png_path.name}\nRun: python scripts/heygen/ai_engine_excalidraw_diagrams.py",
            size=16,
            color=ORANGE,
        )
    return s


def append_excalidraw_diagram_slides(prs):
    """Generate PNG/Excalidraw assets and append diagram appendix slides."""
    generate_all(DIAGRAM_DIR)

    s = blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = TEAL
    bg.line.fill.background()
    add_text_box(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(0.8), "Excalidraw diagrams", size=36, bold=True, color=WHITE)
    add_text_box(
        s,
        Inches(0.8),
        Inches(3.1),
        Inches(11.5),
        Inches(1.5),
        "Simplified workflow views — one scenario per slide\n"
        "Editable source: videos/generated/diagrams/*.excalidraw",
        size=16,
        color=TEAL_LT,
    )

    for title, subtitle, file_id in DIAGRAM_SLIDE_META:
        add_diagram_slide(prs, title, subtitle, DIAGRAM_DIR / f"{file_id}.png")

    # Connectivity section (platform + engine + Bedrock — was not in one diagram before)
    s = blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    add_text_box(
        s,
        Inches(0.8),
        Inches(2.0),
        Inches(11.5),
        Inches(0.9),
        "Connectivity architecture",
        size=36,
        bold=True,
        color=WHITE,
    )
    add_text_box(
        s,
        Inches(0.8),
        Inches(2.95),
        Inches(11.5),
        Inches(1.6),
        "Host platform (running today) · AI Engine · Amazon Bedrock\n"
        "Inputs/outputs and multi-call flows between Engine and Bedrock",
        size=16,
        color=BLUE_LT,
    )

    for title, subtitle, file_id in CONNECTIVITY_DIAGRAM_SLIDE_META:
        add_diagram_slide(prs, title, subtitle, DIAGRAM_DIR / f"{file_id}.png")


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1 Title
    s = blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PURPLE
    bg.line.fill.background()
    add_text_box(s, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.9), "Gas Station AI Engine", size=42, bold=True, color=WHITE)
    add_text_box(s, Inches(0.8), Inches(2.55), Inches(11.5), Inches(0.7), "Implementation & Operations Deck", size=28, color=PURPLE_LT)
    add_text_box(
        s,
        Inches(0.8),
        Inches(3.4),
        Inches(11.5),
        Inches(1.4),
        "Components · development · deployment · training · forecasting\n"
        "LLM products (Claude, Bedrock, RAG, MCP) · hosting · scale",
        size=16,
        color=RGBColor(0xE8, 0xD5, 0xF0),
    )
    add_text_box(
        s,
        Inches(0.8),
        Inches(5.2),
        Inches(11.5),
        Inches(0.6),
        "Companion to: gas-station-coo-developer-architecture.pptx",
        size=13,
        color=ACCENT_LINE,
    )

    # 2 What is the engine
    s = blank_slide(prs)
    slide_header(s, "What is the AI engine?", "A separate service — not part of the Next.js app")
    add_bullet_box(
        s,
        Inches(0.45),
        Inches(1.35),
        Inches(6.0),
        Inches(5.2),
        "The engine does",
        [
            "Pull data from POS, fuel, payroll systems",
            "Store raw history + daily features (its own database)",
            "Run forecasting and anomaly detection",
            "Compose ranked action cards with $ impact",
            "Use an LLM to explain actions in plain English",
            "Run overnight jobs per tenant / timezone",
            "Write curated results back to the platform REST API",
        ],
        fill=TEAL_LT,
        title_color=TEAL,
    )
    add_bullet_box(
        s,
        Inches(6.75),
        Inches(1.35),
        Inches(6.1),
        Inches(5.2),
        "The engine does NOT",
        [
            "Host the owner dashboard (platform does)",
            "Bill customers via Stripe (platform does)",
            "Store raw POS transactions in platform Postgres",
            "Expose LLM API keys to the browser",
            "Replace Clerk authentication for owners",
            "Become the system of record for stations",
        ],
        fill=ORANGE_LT,
        title_color=ORANGE,
    )

    # 3 Contract with platform
    s = blank_slide(prs)
    slide_header(s, "Contract with the platform", "Read registry · write brief · never swap roles")

    boxes = [
        (0.4, 1.5, 3.8, 2.0, BLUE_LT, BLUE, "READ from platform\n\n• Station list\n• Integration registry\n• Prior recommendation status\n• Owner Accept/Dismiss"),
        (4.7, 1.5, 3.8, 2.0, TEAL_LT, TEAL, "PROCESS in engine\n\n• Ingest & normalize\n• Features & models\n• Compose + explain"),
        (9.0, 1.5, 3.8, 2.0, GREEN_LT, GREEN, "WRITE to platform\n\n• Daily metrics (1/station/day)\n• Recommendation cards\n• Service JWT + X-Tenant-ID"),
    ]
    for left, top, w, h, fill, line, text in boxes:
        sh = add_rounded_rect(s, Inches(left), Inches(top), Inches(w), Inches(h), fill, line=line, line_w=2)
        label_shape(sh, text, size=12, bold=True, color=DARK, align=PP_ALIGN.LEFT)

    add_text_box(
        s,
        Inches(0.45),
        Inches(4.0),
        Inches(12.4),
        Inches(2.5),
        "One-sentence contract\n"
        "The platform tells the engine what to connect to. The engine tells the platform what to do today.\n\n"
        "Repo (proposed): gas-station-ai-engine  ·  Blueprint reference: Gas_Station_AI project",
        size=15,
        color=DARK,
    )

    # 4 Eight components
    s = blank_slide(prs)
    slide_header(s, "Eight engine components", "From connectors to write-back client")

    comps = [
        ("C1 Connector framework", "XL", "Plugin per vendor: API, SFTP, or file upload"),
        ("C2 Raw + feature store", "L", "Immutable raw zone + cleaned daily features"),
        ("C3 Forecasting service", "L", "Fuel demand, category sales, labor curves"),
        ("C4 Pricing + anomalies", "M", "Shrink, cash variance, fuel margin signals"),
        ("C5 Recommendation composer", "M", "Rank actions; chain-level cards (null station)"),
        ("C6 LLM explainer + chat", "L", "Plain-English titles; grounded Q&A"),
        ("C7 Orchestrator", "M", "Nightly cron, idempotent runs, retries"),
        ("C8 Platform client SDK", "S", "JWT write-back + REST wrappers"),
    ]
    for i, (name, effort, body) in enumerate(comps):
        col = i % 4
        row = i // 4
        left = Inches(0.35 + col * 3.2)
        top = Inches(1.35 + row * 2.75)
        fill = PURPLE_LT if row == 0 else TEAL_LT
        line = PURPLE if row == 0 else TEAL
        card = add_rounded_rect(s, left, top, Inches(3.05), Inches(2.5), fill, line=line)
        tf = card.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{name}  [{effort}]"
        set_run(r, size=13, bold=True, color=line)
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = body
        set_run(r, size=12, color=DARK)

    # 5 Connector lifecycle
    s = blank_slide(prs)
    slide_header(s, "Connector lifecycle", "How data enters the engine from each store")

    steps = [
        ("1", "Discover", "Load integrations\nfrom platform REST API"),
        ("2", "Resolve secrets", "credentials_ref →\nSecrets Manager"),
        ("3", "Extract", "API poll, SFTP,\nor CSV file drop"),
        ("4", "Normalize", "Canonical schemas:\nsales, fuel, labor"),
        ("5", "Validate", "Dates, row counts,\nstation match"),
        ("6", "Load", "Append raw store;\nrebuild features"),
        ("7", "Report", "PATCH integration\nsync status on platform"),
    ]
    for i, (num, title, body) in enumerate(steps):
        left = Inches(0.25 + i * 1.85)
        bubble = add_rounded_rect(s, left, Inches(1.9), Inches(1.7), Inches(2.6), WHITE, line=TEAL, line_w=1.5)
        tf = bubble.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        set_run(r, size=18, bold=True, color=TEAL)
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = title
        set_run(r, size=12, bold=True, color=NAVY)
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = body
        set_run(r, size=10, color=GRAY)

    add_text_box(
        s,
        Inches(0.4),
        Inches(5.0),
        Inches(12.4),
        Inches(1.6),
        "MVP connector priority\n"
        "P0 — Verifone Commander (POS) via CSV then API  ·  P0 — Gilbarco Passport (fuel) daily file\n"
        "P1 — When I Work (payroll/staffing)  ·  P2 — lottery, car wash, foodservice",
        size=14,
        color=DARK,
    )

    # 6 Data stores & volume
    s = blank_slide(prs)
    slide_header(s, "Where data lives (volume & retention)", "Engine data store is separate from platform RDS")

    add_bullet_box(
        s,
        Inches(0.4),
        Inches(1.35),
        Inches(6.0),
        Inches(5.2),
        "Storage layout",
        [
            "Raw landing zone: S3 (Parquet) partitioned by tenant / station / date",
            "Feature store: Timescale or Postgres TSDB — optional, engine-owned",
            "Model artifacts: versioned in S3 (model_version pointer)",
            "Platform DB: only curated daily metrics + recommendation rows",
            "No raw receipt lines cross into the platform database",
        ],
        fill=TEAL_LT,
        title_color=TEAL,
    )
    add_bullet_box(
        s,
        Inches(6.75),
        Inches(1.35),
        Inches(6.1),
        Inches(5.2),
        "Volume expectations (planning)",
        [
            "1 store ≈ 1 daily metric row + a handful of recs per night",
            "Chain (50 stores) = 50 parallel station pipelines inside one tenant job",
            "Many tenants = queue workers; one tenantId per job for isolation",
            "File uploads spike at end-of-day — workers scale horizontally",
            "Chat / refresh (phase 3) is low volume vs overnight batch",
        ],
        fill=PURPLE_LT,
        title_color=PURPLE,
    )

    # 7 Forecasting
    s = blank_slide(prs)
    slide_header(s, "Forecasting & anomalies", "Quantitative models — not the LLM")

    add_bullet_box(
        s,
        Inches(0.4),
        Inches(1.35),
        Inches(6.0),
        Inches(5.2),
        "Forecasting service outputs",
        [
            "Fuel gallons and margin outlook",
            "Inside sales by category",
            "Labor hours vs demand (staffing recs)",
            "Drives ORDERING and STAFFING categories",
            "Tech: Prophet / GBM / baselines first",
            "Cold start (<30 days): cohort medians, lower confidence",
        ],
        fill=BLUE_LT,
        title_color=BLUE,
    )
    add_bullet_box(
        s,
        Inches(6.75),
        Inches(1.35),
        Inches(6.1),
        Inches(5.2),
        "Anomaly & pricing signals",
        [
            "Cash shortage / void rate z-scores per station",
            "Shrink and waste pattern flags",
            "Fuel pricing signals (rules + optional competitor feed)",
            "Category: LOSS_PREVENTION, FUEL_PRICING",
            "Feeds composer before LLM explains",
            "Stale vendor file → low-confidence warning on brief",
        ],
        fill=ORANGE_LT,
        title_color=ORANGE,
    )

    # 8 Composer + LLM
    s = blank_slide(prs)
    slide_header(s, "From signals to owner-readable actions", "Composer ranks · LLM explains (no raw receipts in prompts)")

    flow = [
        (0.5, 2.1, "Model outputs\n+ rules", TEAL_LT, TEAL),
        (3.4, 2.1, "Composer\nrank by $\nimpact", PURPLE_LT, PURPLE),
        (6.3, 2.1, "Structured\nJSON payload", BLUE_LT, BLUE),
        (9.2, 2.1, "LLM explainer\nplain English", ORANGE_LT, ORANGE),
    ]
    for left, top, text, fill, line in flow:
        sh = add_rounded_rect(s, Inches(left), top, Inches(2.5), Inches(2.0), fill, line=line, line_w=2)
        label_shape(sh, text, size=13, bold=True, color=DARK)
    for x in (3.05, 5.95, 8.85):
        add_arrow(s, Inches(x), Inches(3.1), Inches(x + 0.32), Inches(3.1), PURPLE)

    add_text_box(
        s,
        Inches(0.45),
        Inches(4.5),
        Inches(12.3),
        Inches(2.2),
        "LLM rules (MVP)\n"
        "• Input = structured metrics + model features only — never full transaction logs or secrets.\n"
        "• Tool use: get_metrics, get_recommendations, run_forecast — no direct SQL on platform DB.\n"
        "• Phase 3 chat: Next.js server action → engine /v1/chat (keys stay in engine env only).",
        size=14,
        color=DARK,
    )

    # 8b Two kinds of “AI” (layman)
    s = blank_slide(prs)
    slide_header(s, "Two kinds of “AI” in this engine", "Plain English — forecasting math vs. chat/explain language model")
    add_bullet_box(
        s,
        Inches(0.4),
        Inches(1.35),
        Inches(6.0),
        Inches(5.2),
        "① Forecasting & anomaly models (we train these)",
        [
            "Purpose: predict fuel sales, spot shrink, rank $ impact",
            "Examples: Prophet, gradient boosting, simple rules",
            "Trained weekly on YOUR store history in AWS",
            "Saved as files in S3 (like a spreadsheet of coefficients)",
            "Runs on our ECS workers — not ChatGPT, not Gemini",
            "This is classic machine learning — numbers in, numbers out",
        ],
        fill=BLUE_LT,
        title_color=BLUE,
    )
    add_bullet_box(
        s,
        Inches(6.75),
        Inches(1.35),
        Inches(6.1),
        Inches(5.2),
        "② LLM — language model (we rent, not train, at MVP)",
        [
            "Purpose: turn ranked actions into plain-English explanations",
            "Examples: Claude, GPT, Gemini — brand names owners may know",
            "Pre-trained by Anthropic / OpenAI / Google — we do not rebuild it",
            "We call it over HTTPS (API) with a secret key",
            "Optional chat: owner asks “why premium pricing?”",
            "Fine-tuning the LLM is optional later — not required for pilot",
        ],
        fill=PURPLE_LT,
        title_color=PURPLE,
    )

    # 8c Which LLM product & model names
    s = blank_slide(prs)
    slide_header(s, "Which LLM product & model name?", "What we actually buy — and what Hugging Face / Llama mean")
    add_text_box(
        s,
        Inches(0.45),
        Inches(1.15),
        Inches(12.3),
        Inches(0.5),
        "MVP default (recommended): ONE call to Amazon Bedrock with modelId = Claude 3.5 Sonnet. Bedrock is the AWS API; Claude is the model running inside it — not two separate calls.",
        size=13,
        color=TEAL,
        bold=True,
    )
    providers = [
        ("AWS Bedrock", "AWS’s LLM menu", "Claude, Llama, Mistral via one AWS bill\nNo GPU servers to manage\nKeeps calls inside your AWS account"),
        ("Anthropic API", "Direct from Anthropic", "claude-sonnet / claude-haiku\nSimple if not using Bedrock yet\nSeparate vendor contract"),
        ("OpenAI API", "ChatGPT family", "gpt-4o / gpt-4o-mini\nStrong tool calling\nCommon default in many products"),
        ("Google Gemini", "Via Vertex AI", "gemini-1.5-pro / flash\nGood if already on GCP\nNot required for gas-station MVP"),
        ("Meta Llama", "Open-weight model", "Llama 3.x on Bedrock OR self-host\nCheaper at huge scale; more ops work\nNot “Llama” the app — it’s the model file"),
        ("Hugging Face", "Model library / hub", "Where open models are published\nWe do NOT host the product on HF for MVP\nOptional: download weights if self-hosting Llama"),
    ]
    for i, (name, tag, body) in enumerate(providers):
        col = i % 3
        row = i // 3
        left = Inches(0.35 + col * 4.25)
        top = Inches(1.75 + row * 2.55)
        card = add_rounded_rect(s, left, top, Inches(4.05), Inches(2.35), WHITE, line=PURPLE if row == 0 else TEAL, line_w=1.5)
        tf = card.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{name}  ({tag})"
        set_run(r, size=12, bold=True, color=PURPLE if row == 0 else TEAL)
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = body
        set_run(r, size=11, color=DARK)

    # 8d How we choose the model
    s = blank_slide(prs)
    slide_header(s, "How we choose the LLM", "Decision checklist — not hype, fit for gas-station briefs")
    criteria = [
        ("Structured output", "Must return clean JSON for explanation field", "Claude Sonnet, GPT-4o — high pass rate"),
        ("Plain English", "Owners are not data scientists", "All top-tier models OK; prompt matters more"),
        ("Cost at scale", "50 stores × nightly brief + some chat", "Start Sonnet; move haiku/mini for bulk; Llama later if $ tight"),
        ("AWS alignment", "Platform already on AWS Amplify + ECS", "Bedrock = one invoice, IAM, no extra data egress"),
        ("Data handling", "No raw receipts in prompts", "Any provider OK if we only send summaries"),
        ("Vendor flexibility", "Avoid lock-in", "LiteLLM adapter in engine — swap model ID in config"),
    ]
    for i, (title, why, pick) in enumerate(criteria):
        top = Inches(1.2 + i * 0.95)
        card = add_rounded_rect(s, Inches(0.45), top, Inches(12.2), Inches(0.85), WHITE if i % 2 else GRAY_LT, line=TEAL, line_w=1)
        tf = card.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{title} — {why}"
        set_run(r, size=13, bold=True, color=NAVY)
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = f"Typical pick: {pick}"
        set_run(r, size=11, color=GRAY)

    # 8e Where things are hosted (layman)
    s = blank_slide(prs)
    slide_header(s, "Where is everything hosted?", "Layman map — your servers vs. rented LLM brain")
    add_bullet_box(
        s,
        Inches(0.4),
        Inches(1.35),
        Inches(6.0),
        Inches(5.2),
        "Runs on OUR AWS account",
        [
            "Engine API + overnight workers → ECS Fargate (containers)",
            "Job queue → SQS; schedule → EventBridge",
            "Store files & model files → S3",
            "Secrets (API keys, vendor passwords) → Secrets Manager",
            "Forecasting training jobs → same ECS / batch (CPU)",
            "Platform dashboard stays separate (Amplify + REST API)",
        ],
        fill=TEAL_LT,
        title_color=TEAL,
    )
    add_bullet_box(
        s,
        Inches(6.75),
        Inches(1.35),
        Inches(6.1),
        Inches(5.2),
        "Rented per API call (not on our disk)",
        [
            "LLM inference → Amazon Bedrock OR vendor API",
            "We pay per 1,000 tokens (words-ish), not per server",
            "Model weights live at Amazon/Anthropic — we never download GPT",
            "Hugging Face is NOT where owners log in — it’s a model catalog",
            "Self-hosting Llama on GPU EC2 is optional Phase 4 cost play",
            "Browser never holds LLM keys — only engine server does",
        ],
        fill=ORANGE_LT,
        title_color=ORANGE,
    )

    # 8f AWS Bedrock & cost (plain English)
    s = blank_slide(prs)
    slide_header(s, "AWS Bedrock & LLM costs", "Yes — AWS already sells LLM access; we use it like electricity")
    add_text_box(
        s,
        Inches(0.45),
        Inches(1.15),
        Inches(12.2),
        Inches(0.9),
        "Amazon Bedrock is AWS’s “LLM store.” You pick a model ID (e.g. anthropic.claude-3-5-sonnet), call an API, AWS meters input + output tokens and bills your account — no GPU cluster to babysit.",
        size=14,
        color=DARK,
    )
    cost_rows = [
        ("Nightly brief (per store)", "Small JSON context + 3-sentence explanation", "Cents per store per night at pilot scale"),
        ("Chain (50 stores)", "50 parallel explain calls in one tenant job", "Predictable; batch off-peak"),
        ("Owner chat", "Bursty, shorter threads", "Lower volume than overnight until adoption grows"),
        ("Forecasting ML", "ECS CPU only — scikit-learn / Prophet", "No Bedrock charge — separate compute $"),
        ("What we watch", "Tokens per tenant, model ID, error rate", "CloudWatch + budget alarms"),
    ]
    for i, (what, how, note) in enumerate(cost_rows):
        top = Inches(2.15 + i * 0.95)
        card = add_rounded_rect(s, Inches(0.45), top, Inches(12.2), Inches(0.82), GREEN_LT if i == 4 else WHITE, line=GREEN if i == 4 else BLUE, line_w=1.5)
        tf = card.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{what}  ·  {how}"
        set_run(r, size=13, bold=True, color=NAVY)
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = note
        set_run(r, size=11, color=GRAY)

    # 8f-b Bedrock vs Claude (plain English)
    s = blank_slide(prs)
    slide_header(s, "Bedrock vs Claude — do we need both?", "No — one API call in MVP; two names for two different layers")
    add_bullet_box(
        s,
        Inches(0.4),
        Inches(1.2),
        Inches(6.0),
        Inches(5.4),
        "Short answer",
        [
            "We do NOT call Bedrock and Claude as two separate paid hops",
            "Our engine makes ONE HTTPS request to Amazon Bedrock",
            "In that request we set modelId = anthropic.claude-3-5-sonnet",
            "Bedrock runs the Claude model inside AWS and returns text/JSON",
            "Claude is the external LLM (made by Anthropic) — we never host its weights",
            "See diagram: 11-bedrock-claude-explained.excalidraw",
        ],
        fill=ORANGE_LT,
        title_color=ORANGE,
    )
    add_bullet_box(
        s,
        Inches(6.75),
        Inches(1.2),
        Inches(6.1),
        Inches(5.4),
        "What each is for",
        [
            "Bedrock → HOW we access: AWS API, IAM, token billing, model catalog",
            "Claude → WHAT writes English: explanation field, chat replies",
            "Prophet/GBM (separate) → forecast numbers — not language",
            "Path B (optional): call api.anthropic.com directly — same Claude, no Bedrock",
            "MVP pick: Bedrock + Claude Sonnet — one AWS bill, one integration",
            "Analogy: Bedrock = app store · Claude = the app you chose",
        ],
        fill=PURPLE_LT,
        title_color=PURPLE,
    )

    # 8f-c Platform · Engine · Bedrock connectivity (summary before diagrams)
    s = blank_slide(prs)
    slide_header(s, "Connectivity: platform · engine · Bedrock", "Three systems — who talks to whom (host platform is live today)")
    add_bullet_box(
        s,
        Inches(0.4),
        Inches(1.2),
        Inches(6.0),
        Inches(5.4),
        "Platform ↔ AI Engine",
        [
            "Platform = Next.js + REST API (mosc-temp) — owner dashboard",
            "Engine READs: stations, integrations, prior recommendation status",
            "Engine WRITEs: daily metrics, recommendation cards, sync status",
            "Platform triggers engine on-demand: POST /v1/chat, /v1/runs/trigger",
            "All via server actions + service JWT — browser never calls engine or Bedrock",
            "See diagrams: 12-platform-engine-bedrock-connectivity",
        ],
        fill=BLUE_LT,
        title_color=BLUE,
    )
    add_bullet_box(
        s,
        Inches(6.75),
        Inches(1.2),
        Inches(6.1),
        Inches(5.4),
        "AI Engine ↔ Amazon Bedrock",
        [
            "Only the engine holds Bedrock credentials — not the platform",
            "OUT: structured JSON prompt per station explanation (×N nightly)",
            "IN: explanation text / JSON from Claude inside Bedrock",
            "Chat may loop 2–4 Bedrock calls (tool use + RAG + final answer)",
            "Platform is not in the Bedrock path — engine aggregates then write-back",
            "See diagram: 13-engine-bedrock-io-multicall",
        ],
        fill=ORANGE_LT,
        title_color=ORANGE,
    )

    # 8g RAG explained
    s = blank_slide(prs)
    slide_header(s, "RAG — memory outside the model", "Retrieval is external; the LLM only reads a briefing packet we assemble")
    flow_rag = [
        (0.35, 2.0, "Store data\n(S3 + features)", TEAL_LT, TEAL),
        (3.2, 2.0, "Retriever\npicks today’s\nmetrics + recs", PURPLE_LT, PURPLE),
        (6.05, 2.0, "Context packet\n(JSON summary,\nno raw receipts)", BLUE_LT, BLUE),
        (8.9, 2.0, "Bedrock API\n(Claude inside)", ORANGE_LT, ORANGE),
        (11.75, 2.0, "Owner text\nexplanation", GREEN_LT, GREEN),
    ]
    for left, top, text, fill, line in flow_rag:
        sh = add_rounded_rect(s, Inches(left), top, Inches(2.55), Inches(2.15), fill, line=line, line_w=2)
        label_shape(sh, text, size=11, bold=True, color=DARK)
    for x in (2.95, 5.8, 8.65, 11.5):
        add_arrow(s, Inches(x), Inches(3.05), Inches(x + 0.22), Inches(3.05), PURPLE)
    add_bullet_box(
        s,
        Inches(0.4),
        Inches(4.55),
        Inches(12.3),
        Inches(2.1),
        "Key ideas (non-technical)",
        [
            "RAG does NOT train the model inside weights — it fetches facts first, then asks the LLM to explain",
            "Built into our engine (Python), not inside Gemini/Claude itself",
            "Vector DB (e.g. OpenSearch) is optional Phase 3+ for long chat history — MVP uses structured JSON RAG",
            "Fine-tuning = actually changing model weights — expensive, later, only if prompts+RAG insufficient",
        ],
        fill=GRAY_LT,
        title_color=GRAY,
    )

    # 8h MCP & tools
    s = blank_slide(prs)
    slide_header(s, "Tools, function calling & MCP", "How the agent “does things” beyond typing — integration pattern")
    add_bullet_box(
        s,
        Inches(0.4),
        Inches(1.35),
        Inches(6.0),
        Inches(5.2),
        "MVP — function calling (built-in)",
        [
            "LLM returns a tool name + arguments (JSON)",
            "Engine runs Python functions: get_metrics, run_forecast",
            "Results fed back to LLM for final answer",
            "Hardcoded in gas-station-ai-engine repo",
            "Same pattern OpenAI / Anthropic / Bedrock support",
            "No owner browser access — server-side only",
        ],
        fill=BLUE_LT,
        title_color=BLUE,
    )
    add_bullet_box(
        s,
        Inches(6.75),
        Inches(1.35),
        Inches(6.1),
        Inches(5.2),
        "Later — MCP (Model Context Protocol)",
        [
            "Open standard to plug in new data sources as “tools”",
            "Example: competitor price feed, weather, lottery API",
            "Each connector can expose an MCP server",
            "Engine agent discovers tools without rewriting core chat",
            "MCP is NOT the LLM — it’s how we attach capabilities",
            "RAG = what to read · MCP/tools = what to invoke",
        ],
        fill=PURPLE_LT,
        title_color=PURPLE,
    )

    # 8i Programming stack & effort
    s = blank_slide(prs)
    slide_header(s, "Programming stack & effort", "What engineers build — languages, libraries, person-weeks")
    add_bullet_box(
        s,
        Inches(0.4),
        Inches(1.35),
        Inches(6.0),
        Inches(5.2),
        "Tech stack",
        [
            "Language: Python 3.11+ (connectors, ML, LLM agent)",
            "API: FastAPI — /v1/chat, /v1/runs/trigger, /health",
            "ML: pandas, scikit-learn, Prophet (or similar)",
            "LLM client: boto3 Bedrock OR anthropic/openai SDK",
            "Optional: LiteLLM (one interface, many providers)",
            "Optional: LangChain for agent loops — not required MVP",
            "Platform SDK: Python client for REST write-back",
        ],
        fill=TEAL_LT,
        title_color=TEAL,
    )
    add_bullet_box(
        s,
        Inches(6.75),
        Inches(1.35),
        Inches(6.1),
        Inches(5.2),
        "Indicative effort (LLM slice)",
        [
            "Phase 2: explainer prompt + JSON schema — ~1 week",
            "Phase 2: wire Bedrock + secrets + logging — ~3–5 days",
            "Phase 3: chat endpoint + tool loop — ~2 weeks",
            "Phase 3: RAG retriever over metrics/recs — ~1 week",
            "Phase 4: MCP connector plugin — ~1–2 weeks each",
            "Ongoing: prompt tuning, cost monitoring — part-time",
            "Total engine (all components): ~10–15 person-weeks",
        ],
        fill=GREEN_LT,
        title_color=GREEN,
    )

    # 9 Nightly orchestration
    s = blank_slide(prs)
    slide_header(s, "Overnight orchestration", "One scheduled run per tenant — then per station inside the job")

    layers = [
        "EventBridge / cron lists tenants with gas module enabled",
        "Respect each station timezone + gas_daily_brief_hour_local",
        "Enqueue tenant job → worker pulls integrations for all stations",
        "Per station: extract → normalize → features → models → compose → LLM → write-back",
        "Partial failure isolation: one bad CSV does not kill the whole chain",
        "Idempotent: same calendar day rerun does not duplicate metrics",
        "Log source_model_run_id on every recommendation for traceability",
    ]
    for i, line in enumerate(layers):
        top = Inches(1.25 + i * 0.78)
        card = add_rounded_rect(s, Inches(0.5), top, Inches(12.3), Inches(0.68), WHITE, line=TEAL, line_w=1.5)
        label_shape(card, f"{i + 1}.  {line}", size=13, bold=False, color=DARK, align=PP_ALIGN.LEFT)

    # 10 Training pipeline
    s = blank_slide(prs)
    slide_header(s, "Training & model refresh", "Weekly batch training — nightly jobs only infer")

    add_text_box(
        s,
        Inches(0.45),
        Inches(1.2),
        Inches(12.3),
        Inches(0.55),
        "Important: “Training” = quantitative models on store history. LLM uses prompting + RAG first; fine-tuning is optional in phase 3+.",
        size=13,
        color=ORANGE,
        bold=True,
    )

    train_steps = [
        ("Extract", "Features for stations with ≥30 days history"),
        ("Train", "Fuel, category sales, labor artifacts → S3 versioned"),
        ("Evaluate", "Holdout last 7 days; keep old model if MAPE degrades"),
        ("Promote", "Production pointer for nightly orchestrator"),
        ("Log", "source_model_run_id on each recommendation row"),
    ]
    for i, (title, body) in enumerate(train_steps):
        left = Inches(0.4 + i * 2.55)
        card = add_rounded_rect(s, left, Inches(2.0), Inches(2.35), Inches(2.8), TEAL_LT, line=TEAL)
        tf = card.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = title
        set_run(r, size=16, bold=True, color=TEAL)
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = body
        set_run(r, size=12, color=DARK)

    add_bullet_box(
        s,
        Inches(0.4),
        Inches(5.1),
        Inches(12.3),
        Inches(1.7),
        "Learning from owners (continuous improvement)",
        [
            "ACCEPTED / DONE → positive label for recommendation type",
            "DISMISSED + feedback → down-rank similar future actions",
            "Actual vs expected profit next day → calibrate $ impact estimates",
        ],
        fill=GREEN_LT,
        title_color=GREEN,
    )

    # 11 Deployment architecture
    s = blank_slide(prs)
    slide_header(s, "Deployment & hosting", "Separate stack from the platform — same cloud patterns")

    deploy = [
        ("Engine API + workers", "ECS Fargate or Lambda + SQS/RQ queue\n(match event-site-manager-batch-jobs)"),
        ("Scheduler", "EventBridge rules per timezone window"),
        ("Raw + features", "S3 + optional Timescale (not platform RDS)"),
        ("LLM provider keys", "Engine environment only — never in mosc-temp"),
        ("Secrets", "AWS Secrets Manager for vendor + platform JWT"),
        ("IaC", "application-deployments → gas-station-ai-engine stack"),
    ]
    for i, (title, body) in enumerate(deploy):
        col = i % 3
        row = i // 3
        left = Inches(0.4 + col * 4.25)
        top = Inches(1.35 + row * 2.75)
        card = add_rounded_rect(s, left, top, Inches(4.05), Inches(2.45), WHITE, line=PURPLE, line_w=2)
        tf = card.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        set_run(r, size=15, bold=True, color=PURPLE)
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = body
        set_run(r, size=12, color=DARK)

    # 12 Scale & multi-tenant
    s = blank_slide(prs)
    slide_header(s, "Scale, traffic & many stores", "How the engine handles growth")

    add_bullet_box(
        s,
        Inches(0.4),
        Inches(1.35),
        Inches(6.0),
        Inches(5.2),
        "Horizontal scale",
        [
            "Queue-backed workers add capacity for more tenants",
            "Station pipelines run in parallel within a tenant job",
            "Connector failures isolated per integration row",
            "Dead-letter queue + alerts for ops (phase 4)",
            "Write-back retries with JWT refresh (401 → new token)",
            "Nightly batch dominates cost; chat is bursty but smaller",
        ],
        fill=TEAL_LT,
        title_color=TEAL,
    )
    add_bullet_box(
        s,
        Inches(6.75),
        Inches(1.35),
        Inches(6.1),
        Inches(5.2),
        "Multi-tenant isolation",
        [
            "One tenantId per orchestration job — no cross-tenant mixing",
            "S3 prefixes: tenant / station / date partitions",
            "Service JWT scoped to gas-station write endpoints",
            "X-Tenant-ID on every platform REST call",
            "Chain view: engine emits station_id=null for cross-store recs",
            "Per-station timezone drives when brief is due",
        ],
        fill=PURPLE_LT,
        title_color=PURPLE,
    )

    # 13 On-demand APIs
    s = blank_slide(prs)
    slide_header(s, "On-demand: refresh & chat (phase 3)", "Platform server actions call engine — not the browser")

    rows = [
        ("Refresh brief button", "Server action → POST /v1/runs/trigger", "Re-run pipeline for tenant or one station"),
        ("AI Manager chat", "Server action → POST /v1/chat", "Grounded Q&A with tool use"),
        ("Engine reads context", "Platform SDK + service JWT", "Metrics, recs, integration status"),
        ("Config per tenant", "tenant_settings.gas_ai_engine_*", "base_url, api_key_ref, brief hour"),
    ]
    for i, (action, path, note) in enumerate(rows):
        top = Inches(1.35 + i * 1.35)
        card = add_rounded_rect(s, Inches(0.5), top, Inches(12.3), Inches(1.15), WHITE, line=BLUE, line_w=1.5)
        tf = card.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{action}  —  {path}"
        set_run(r, size=14, bold=True, color=NAVY)
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = note
        set_run(r, size=12, color=GRAY)

    # 14 Security
    s = blank_slide(prs)
    slide_header(s, "Security essentials", "Financial data — treat as production-grade from day one")

    rules = [
        ("Tenant isolation", "Jobs and object-store paths scoped by tenantId"),
        ("Credentials", "Only credentials_ref in platform DB; engine uses Secrets Manager"),
        ("JWT rotation", "Service account for write-back; same rotation as other API users"),
        ("LLM payload", "Summaries + structured features — no PII, no secrets"),
        ("Owner auth", "Clerk required before real client financial data (feasibility risk)"),
        ("Audit", "source_model_run_id links every card to a model run"),
    ]
    for i, (title, body) in enumerate(rules):
        col = i % 2
        row = i // 2
        left = Inches(0.45 + col * 6.4)
        top = Inches(1.35 + row * 1.85)
        card = add_rounded_rect(s, left, top, Inches(6.05), Inches(1.6), GRAY_LT, line=GRAY)
        tf = card.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        set_run(r, size=14, bold=True, color=NAVY)
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = body
        set_run(r, size=12, color=DARK)

    # 15 Phases
    s = blank_slide(prs)
    slide_header(s, "Implementation phases", "~10–15 person-weeks to pilot-quality brief")

    phases = [
        ("0 Foundation", "1–2 wk", "Repo, platform SDK, JWT smoke test, one manual write-back"),
        ("1 MVP ingest", "3–4 wk", "POS + fuel connectors, feature store, rule-based recs"),
        ("2 ML brief", "4–6 wk", "Forecast + anomaly + LLM explain + overnight orchestrator"),
        ("3 Chat + refresh", "2–3 wk", "/v1/chat, server actions, grounded tool use"),
        ("4 Hardening", "4+ wk", "Multi-tenant scale, monitoring, 2nd vendor, feedback loop"),
    ]
    for i, (name, dur, scope) in enumerate(phases):
        top = Inches(1.25 + i * 1.15)
        card = add_rounded_rect(s, Inches(0.5), top, Inches(12.3), Inches(1.0), WHITE, line=PURPLE, line_w=2)
        tf = card.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"Phase {name}  ({dur})"
        set_run(r, size=14, bold=True, color=PURPLE)
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = scope
        set_run(r, size=12, color=DARK)

    # 16 Acceptance
    s = blank_slide(prs)
    slide_header(s, "Acceptance criteria (Loop 4 complete)", "Engine is “done enough” when the owner sees a real brief")

    criteria = [
        "Overnight job runs for demo_gas_station_001 without manual SQL",
        "Daily brief shows engine metrics + ≥3 recommendations with LLM explanation",
        "Integration rows show last_sync_status = SUCCESS after connector run",
        "Owner Accept/Dismiss updates platform; next run reads new statuses",
        "All engine→platform calls use service JWT + X-Tenant-ID",
        "No secrets in git; LLM keys only in engine environment",
    ]
    for i, c in enumerate(criteria):
        top = Inches(1.3 + i * 0.88)
        card = add_rounded_rect(s, Inches(0.55), top, Inches(12.1), Inches(0.75), GREEN_LT if i % 2 == 0 else WHITE, line=GREEN)
        label_shape(card, f"✓  {c}", size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT)

    # Diagram appendix (Excalidraw + embedded PNGs)
    append_excalidraw_diagram_slides(prs)

    # Close
    s = blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PURPLE
    bg.line.fill.background()
    add_text_box(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.9), "Ingest · model · explain · write back", size=30, bold=True, color=WHITE)
    add_text_box(
        s,
        Inches(0.8),
        Inches(3.1),
        Inches(11.5),
        Inches(2.2),
        "Primary references\n"
        "• gas_station_ai_engine_prd.html\n"
        "• gas_station_ai_engine_workflow.excalidraw\n"
        "• Simplified diagrams: videos/generated/diagrams/*.excalidraw\n"
        "• LLM: Bedrock / Claude · RAG external · MCP for tools (Phase 4+)\n"
        "• Companion deck: gas-station-coo-developer-architecture.pptx",
        size=14,
        color=PURPLE_LT,
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(OUT))
        print(f"Wrote {OUT}")
        return OUT
    except PermissionError:
        alt = OUT.with_name("gas-station-coo-ai-engine-implementation-with-diagrams.pptx")
        prs.save(str(alt))
        print(f"Wrote {alt} (original locked)")
        return alt


if __name__ == "__main__":
    build()
