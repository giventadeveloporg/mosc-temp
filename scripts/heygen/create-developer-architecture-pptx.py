#!/usr/bin/env python3
"""
Gas Station COO — Developer Architecture Infographic Deck

Builds a high-level PowerPoint for engineers from the same material as the
developer walkthrough video (gas_station_video_scripts_developer.html).

Focus: how AI engine ↔ platform backend ↔ frontend interact, how data moves,
and what each major component does — in plain English with Excalidraw-style diagrams.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ---------- palette (developer / docs blue, not warm-cream AI cliché) ----------
NAVY = RGBColor(0x0D, 0x3B, 0x66)
BLUE = RGBColor(0x00, 0x66, 0xCC)
BLUE_LT = RGBColor(0xE8, 0xF4, 0xF8)
TEAL = RGBColor(0x0D, 0x94, 0x88)
TEAL_LT = RGBColor(0xE0, 0xF7, 0xF4)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
ORANGE_LT = RGBColor(0xFE, 0xF5, 0xE7)
PURPLE = RGBColor(0x6C, 0x34, 0x8D)
PURPLE_LT = RGBColor(0xF5, 0xEE, 0xF8)
GREEN = RGBColor(0x1E, 0x84, 0x4E)
GREEN_LT = RGBColor(0xE8, 0xF8, 0xEF)
GRAY = RGBColor(0x4A, 0x55, 0x68)
GRAY_LT = RGBColor(0xF4, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x23, 0x32)
ACCENT_LINE = RGBColor(0x94, 0xA3, 0xB8)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUT = Path(__file__).resolve().parents[2] / (
    "documentation/tenant_management/gas_station_site/videos/generated/decks/"
    "gas-station-coo-developer-architecture.pptx"
)


def set_run(run, size=18, bold=False, color=DARK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_text_box(slide, left, top, width, height, text, *, size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box


def add_rounded_rect(slide, left, top, width, height, fill, *, line=None, line_w=1.25):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    # Soften corner
    try:
        shape.adjustments[0] = 0.15
    except Exception:
        pass
    return shape


def label_shape(shape, text, *, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.get_or_add_bodyPr().set(qn("a:anchor"), "ctr" if valign == MSO_ANCHOR.MIDDLE else "t")
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)


def add_bullet_box(slide, left, top, width, height, title, bullets, *, fill=BLUE_LT, title_color=NAVY, body_color=DARK):
    card = add_rounded_rect(slide, left, top, width, height, fill, line=BLUE, line_w=1)
    tf = card.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_run(run, size=15, bold=True, color=title_color)
    for b in bullets:
        p = tf.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = f"• {b}"
        set_run(run, size=13, bold=False, color=body_color)
    return card


def add_arrow(slide, x1, y1, x2, y2, color=BLUE):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(2.25)
    return connector


def slide_header(slide, title, subtitle=None):
    # top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    add_text_box(slide, Inches(0.4), Inches(0.18), Inches(12.2), Inches(0.45), title, size=26, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, Inches(0.4), Inches(0.55), Inches(12.2), Inches(0.35), subtitle, size=13, color=BLUE_LT)
    # footer
    add_text_box(
        slide,
        Inches(0.4),
        Inches(7.15),
        Inches(12.4),
        Inches(0.3),
        "Gas Station COO · Developer architecture deck · Derived from developer walkthrough script",
        size=10,
        color=ACCENT_LINE,
    )


def blank_slide(prs):
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ========== 1 Title ==========
    s = blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    add_text_box(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(1), "Gas Station COO", size=44, bold=True, color=WHITE)
    add_text_box(
        s,
        Inches(0.8),
        Inches(2.7),
        Inches(11.5),
        Inches(0.7),
        "Developer Architecture Infographic Deck",
        size=28,
        bold=False,
        color=BLUE_LT,
    )
    add_text_box(
        s,
        Inches(0.8),
        Inches(3.6),
        Inches(11.5),
        Inches(1.2),
        "How the front end, platform backend, and AI engine work together\n"
        "High-level diagrams in plain English — same material as the developer walkthrough video",
        size=16,
        color=RGBColor(0xC5, 0xD8, 0xEA),
    )
    add_text_box(
        s,
        Inches(0.8),
        Inches(5.4),
        Inches(11.5),
        Inches(0.8),
        "Audience: engineers & technical leads  ·  Level: architecture overview (not minute implementation detail)",
        size=13,
        color=ACCENT_LINE,
    )

    # ========== 2 Agenda ==========
    s = blank_slide(prs)
    slide_header(s, "What this deck covers", "A visual map of the system — not an API reference")
    agenda = [
        ("1", "Big picture", "Platform vs AI engine — who owns what"),
        ("2", "Data journey", "Store systems → engine → curated results → dashboard"),
        ("3", "Front end ↔ backend", "How Next.js talks to REST APIs safely"),
        ("4", "Engine ↔ platform", "How write-back and feedback work"),
        ("5", "Building blocks", "Tables, connectors, models, nightly job"),
        ("6", "Developer loop", "Local demo without engine; phases to pilot"),
    ]
    for i, (num, title, body) in enumerate(agenda):
        col = i % 3
        row = i // 3
        left = Inches(0.5 + col * 4.2)
        top = Inches(1.3 + row * 2.7)
        card = add_rounded_rect(s, left, top, Inches(3.9), Inches(2.3), BLUE_LT, line=BLUE)
        tf = card.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = num
        set_run(r, size=28, bold=True, color=BLUE)
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = title
        set_run(r, size=18, bold=True, color=NAVY)
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = body
        set_run(r, size=13, color=GRAY)

    # ========== 3 Big picture split ==========
    s = blank_slide(prs)
    slide_header(s, "Big picture: two systems, one contract", "Raw POS data never lands in the platform database")

    left = add_rounded_rect(s, Inches(0.4), Inches(1.3), Inches(6.0), Inches(5.3), BLUE_LT, line=BLUE)
    right = add_rounded_rect(s, Inches(6.9), Inches(1.3), Inches(6.0), Inches(5.3), TEAL_LT, line=TEAL)

    add_text_box(s, Inches(0.65), Inches(1.45), Inches(5.5), Inches(0.4), "THIS PLATFORM (mosc-temp)", size=16, bold=True, color=NAVY)
    add_text_box(
        s,
        Inches(0.65),
        Inches(1.95),
        Inches(5.5),
        Inches(4.3),
        "Owns:\n"
        "• Postgres tables (stations, metrics, recommendations)\n"
        "• REST APIs with tenant isolation\n"
        "• Next.js admin UI (daily brief, compare, billing)\n"
        "• Stripe subscription per location\n"
        "• Owner actions: Accept / Dismiss / Done + feedback\n\n"
        "Think of it as:\n"
        "the control room + the books + the customer portal",
        size=14,
        color=DARK,
    )

    add_text_box(s, Inches(7.15), Inches(1.45), Inches(5.5), Inches(0.4), "AI ENGINE (separate service)", size=16, bold=True, color=TEAL)
    add_text_box(
        s,
        Inches(7.15),
        Inches(1.95),
        Inches(5.5),
        Inches(4.3),
        "Owns:\n"
        "• Connectors to POS / fuel / payroll systems\n"
        "• Feature store + forecasting models\n"
        "• Anomaly detection + recommendation composer\n"
        "• LLM that writes plain-English action cards\n"
        "• Nightly orchestration per tenant timezone\n\n"
        "Think of it as:\n"
        "the sensors + the brains that write a short brief",
        size=14,
        color=DARK,
    )

    # middle contract badge
    badge = add_rounded_rect(s, Inches(5.35), Inches(3.5), Inches(2.6), Inches(1.2), WHITE, line=ORANGE, line_w=2)
    label_shape(badge, "ONLY COUPLING\nTyped write-back\nDTOs + service JWT", size=11, bold=True, color=ORANGE)

    # ========== 4 Architecture flow ==========
    s = blank_slide(prs)
    slide_header(s, "Architecture at a glance", "Excalidraw-style flow: store systems → AI engine → platform → owner")

    boxes = [
        (0.35, 2.4, 2.4, 2.2, ORANGE_LT, ORANGE, "Store systems\n\nPOS · Fuel · Payroll\n(Verifone, Gilbarco,\nWhen I Work, …)"),
        (3.15, 2.4, 2.7, 2.2, TEAL_LT, TEAL, "AI Engine\n\nIngest → Normalize\nFeatures → Models\nCompose → Explain"),
        (6.25, 2.4, 2.7, 2.2, BLUE_LT, BLUE, "Platform API\n\nREST API\ntenant-scoped\nJWT + X-Tenant-ID"),
        (9.4, 2.4, 3.4, 2.2, GREEN_LT, GREEN, "Admin UI / Owner\n\nDaily brief\nAccept / Dismiss\nCompare & billing"),
    ]
    for left, top, w, h, fill, line, text in boxes:
        sh = add_rounded_rect(s, Inches(left), Inches(top), Inches(w), Inches(h), fill, line=line, line_w=2)
        label_shape(sh, text, size=13, bold=True, color=DARK)

    # arrows between
    for x in (2.8, 5.9, 9.05):
        add_arrow(s, Inches(x), Inches(3.5), Inches(x + 0.3), Inches(3.5), BLUE)

    add_text_box(
        s,
        Inches(0.4),
        Inches(5.0),
        Inches(12.4),
        Inches(1.6),
        "Rule of thumb\n"
        "• Engine keeps the noisy raw data.\n"
        "• Platform only stores curated daily numbers + a short list of actions.\n"
        "• Owner decisions stay on the platform and later feed the engine’s learning loop.",
        size=14,
        color=DARK,
    )

    # ========== 5 Data collected ==========
    s = blank_slide(prs)
    slide_header(s, "What data is collected (curated, not raw)", "The daily metric is one clean snapshot per station per calendar day")

    items = [
        (BLUE_LT, BLUE, "Fuel", "Gallons, revenue,\nmargin signals"),
        (TEAL_LT, TEAL, "Inside store", "Sales by category,\ntransaction counts"),
        (ORANGE_LT, ORANGE, "Labor", "Hours worked,\nlabor cost estimates"),
        (PURPLE_LT, PURPLE, "Health / risk", "Shrink, waste,\nprofit vs expected"),
        (GREEN_LT, GREEN, "Actions", "Ranked cards with\n$ impact + confidence"),
        (GRAY_LT, GRAY, "Integrations", "Which vendors, sync\nstatus, secret refs"),
    ]
    for i, (fill, line, title, body) in enumerate(items):
        col = i % 3
        row = i // 3
        left = Inches(0.5 + col * 4.2)
        top = Inches(1.35 + row * 2.65)
        card = add_rounded_rect(s, left, top, Inches(3.9), Inches(2.35), fill, line=line, line_w=1.5)
        tf = card.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        set_run(r, size=18, bold=True, color=line)
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = body
        set_run(r, size=14, color=DARK)

    # ========== 6 Frontend ↔ Backend ==========
    s = blank_slide(prs)
    slide_header(s, "Front end ↔ backend", "Browser never talks to the REST API with secrets")

    layers = [
        (0.4, 1.4, 12.5, 1.15, ORANGE_LT, ORANGE, "① Browser (React / Next.js client)", "Shows screens. Handles clicks. Does not call protected APIs directly."),
        (0.4, 2.75, 12.5, 1.15, BLUE_LT, BLUE, "② Server actions / ApiServerActions.ts", "Runs on the server. Adds Authorization + X-Tenant-ID. Uses proxy routes."),
        (0.4, 4.1, 12.5, 1.15, TEAL_LT, TEAL, "③ Next.js proxy  /api/proxy/gas-station-*", "Forwards to backend with JWT retry. Injects tenant context reliably."),
        (0.4, 5.45, 12.5, 1.15, PURPLE_LT, PURPLE, "④ Platform REST API", "CRUD for stations & integrations; brief; compare; PATCH recommendation status."),
    ]
    for left, top, w, h, fill, line, title, body in layers:
        card = add_rounded_rect(s, Inches(left), Inches(top), Inches(w), Inches(h), fill, line=line, line_w=1.5)
        tf = card.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        set_run(r, size=15, bold=True, color=line)
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = body
        set_run(r, size=13, color=DARK)

    # ========== 7 Admin UI map ==========
    s = blank_slide(prs)
    slide_header(s, "Admin UI map", "Where developers will work in the Next.js app")

    ui = [
        ("/admin/gas-station", "Daily brief", "Date + station picker, metric cards, ranked actions", BLUE),
        ("…/stations", "Stations", "List/create/edit store locations under the tenant", TEAL),
        ("…/integrations", "Integrations", "Register POS / fuel / payroll connectors + status", ORANGE),
        ("…/compare", "Compare", "Rank stations for a date range (volume, variance)", PURPLE),
        ("…/billing", "Billing", "Included locations, tier preview, Stripe portal", GREEN),
        ("Feature flag", "Gating", "Only shows when enable_gas_station_module is true", GRAY),
    ]
    for i, (path, title, body, color) in enumerate(ui):
        col = i % 3
        row = i // 3
        left = Inches(0.45 + col * 4.25)
        top = Inches(1.35 + row * 2.7)
        card = add_rounded_rect(s, left, top, Inches(4.0), Inches(2.4), WHITE, line=color, line_w=2)
        tf = card.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        set_run(r, size=18, bold=True, color=color)
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = path
        set_run(r, size=12, bold=True, color=NAVY)
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = body
        set_run(r, size=13, color=DARK)

    # ========== 8 Engine ↔ Platform contract ==========
    s = blank_slide(prs)
    slide_header(s, "AI engine ↔ platform (the handshake)", "Service JWT — not a signed-in human session")

    # left engine, right platform
    eng = add_rounded_rect(s, Inches(0.4), Inches(1.35), Inches(5.8), Inches(3.4), TEAL_LT, line=TEAL, line_w=2)
    label_shape(
        eng,
        "AI ENGINE writes\n\n"
        "POST daily metrics\n"
        "(one row / station / day)\n\n"
        "POST recommendations\n"
        "(action cards + confidence)\n\n"
        "Headers:\nAuthorization: Bearer <service JWT>\nX-Tenant-ID: <tenant>",
        size=13,
        bold=False,
        color=DARK,
        align=PP_ALIGN.LEFT,
    )

    plat = add_rounded_rect(s, Inches(7.1), Inches(1.35), Inches(5.8), Inches(3.4), BLUE_LT, line=BLUE, line_w=2)
    label_shape(
        plat,
        "PLATFORM reads + stores\n\n"
        "Tenant-isolated tables\n"
        "Admin UI shows the brief\n\n"
        "Owner PATCH status\n"
        "PENDING → ACCEPTED /\nDISMISSED / DONE\n\n"
        "Optional feedback_text\n(human-in-the-loop signal)",
        size=13,
        bold=False,
        color=DARK,
        align=PP_ALIGN.LEFT,
    )

    add_text_box(
        s,
        Inches(0.4),
        Inches(5.0),
        Inches(12.5),
        Inches(1.6),
        "Why this design\n"
        "• Keeps secrets and model code out of the web app.\n"
        "• Lets the UI ship and demo with seed data before any engine exists.\n"
        "• Makes the feedback loop visible: owners click Accept → engine can learn what worked.",
        size=14,
        color=DARK,
    )

    # ========== 9 Nightly pipeline ==========
    s = blank_slide(prs)
    slide_header(s, "Nightly job — how a brief is built", "One tenant, then each of its stations (idempotent daily run)")

    steps = [
        ("1", "Pull", "Load active\nintegrations"),
        ("2", "Fetch", "API or file\nfrom vendors"),
        ("3", "Normalize", "Canonical\nschemas"),
        ("4", "Features", "Store daily\nfeature rows"),
        ("5", "Models", "Forecast +\nanomalies"),
        ("6", "Compose", "Rank actions\nby $ impact"),
        ("7", "Explain", "LLM titles\nin plain English"),
        ("8", "Write-back", "POST metrics\n& recs"),
    ]
    for i, (num, title, body) in enumerate(steps):
        left = Inches(0.3 + i * 1.6)
        top = Inches(2.0)
        bubble = add_rounded_rect(s, left, top, Inches(1.45), Inches(2.5), WHITE, line=BLUE, line_w=1.5)
        tf = bubble.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        set_run(r, size=20, bold=True, color=BLUE)
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = title
        set_run(r, size=13, bold=True, color=NAVY)
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = body
        set_run(r, size=11, color=GRAY)

    add_text_box(
        s,
        Inches(0.4),
        Inches(5.0),
        Inches(12.5),
        Inches(1.6),
        "Scheduling\n"
        "Cron checks tenants with the gas module enabled, respects each store’s timezone and preferred brief hour,\n"
        "and isolates failures per station so one bad CSV does not kill the whole chain run.",
        size=14,
        color=DARK,
    )

    # ========== 10 Components C1-C8 ==========
    s = blank_slide(prs)
    slide_header(s, "AI engine components (plain English)", "Eight pieces — from connectors to write-back")

    comps = [
        ("C1 Connectors", "Plug-ins that talk to each vendor (API or file drop)"),
        ("C2 Ingest", "Pull / receive files and turn them into clean rows"),
        ("C3 Feature store", "History the models need (per tenant + station + day)"),
        ("C4 Forecasting", "Guess tomorrow’s fuel / inside sales"),
        ("C5 Anomalies", "Flag weird shrink, cash, or labor patterns"),
        ("C6 Composer", "Turn signals into ranked action cards"),
        ("C7 LLM explainer", "Write the title/body humans read (no raw receipts)"),
        ("C8 Write-back", "Small client that POSTs results to the platform"),
    ]
    for i, (title, body) in enumerate(comps):
        col = i % 4
        row = i // 4
        left = Inches(0.35 + col * 3.2)
        top = Inches(1.35 + row * 2.7)
        card = add_rounded_rect(s, left, top, Inches(3.05), Inches(2.4), BLUE_LT if row == 0 else TEAL_LT, line=BLUE if row == 0 else TEAL)
        tf = card.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        set_run(r, size=15, bold=True, color=NAVY)
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = body
        set_run(r, size=13, color=DARK)

    # ========== 11 Schema ==========
    s = blank_slide(prs)
    slide_header(s, "Platform data model (4 tables)", "Tiny surface area on purpose")

    tables = [
        ("gas_station_location", "Each physical site under a tenant\n(timezone, included_in_subscription)"),
        ("gas_station_integration", "How we connect to a vendor\n(type, mode, secret_ref, sync status)"),
        ("gas_station_daily_metric", "One curated day per station\n(fuel, sales, labor, profit)"),
        ("gas_station_recommendation", "Action cards for owners\n(null station_id = chain-level)"),
    ]
    for i, (name, body) in enumerate(tables):
        left = Inches(0.45 + (i % 2) * 6.4)
        top = Inches(1.35 + (i // 2) * 2.7)
        card = add_rounded_rect(s, left, top, Inches(6.05), Inches(2.4), WHITE, line=NAVY, line_w=2)
        tf = card.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = name
        set_run(r, size=16, bold=True, color=BLUE)
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = body
        set_run(r, size=14, color=DARK)

    # ========== 12 Connectors ==========
    s = blank_slide(prs)
    slide_header(s, "Pilot connectors (priority order)", "Start with files if APIs are delayed — plugins stay the same")

    vendors = [
        ("P0", "Verifone Commander", "POS / inside sales", "CSV first, then REST if available", ORANGE),
        ("P0", "Gilbarco Passport", "Fuel gallons & margin inputs", "Daily file drop (S3 / SFTP)", TEAL),
        ("P1", "When I Work", "Labor hours & staffing", "CSV export after POS+fuel stable", PURPLE),
    ]
    for i, (pri, name, role, note, color) in enumerate(vendors):
        top = Inches(1.35 + i * 1.75)
        card = add_rounded_rect(s, Inches(0.5), top, Inches(12.3), Inches(1.55), WHITE, line=color, line_w=2)
        tf = card.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{pri}  ·  {name}"
        set_run(r, size=18, bold=True, color=color)
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = f"{role}  —  {note}"
        set_run(r, size=14, color=DARK)

    # ========== 13 Feedback loop ==========
    s = blank_slide(prs)
    slide_header(s, "Human-in-the-loop", "Recommendations are suggestions — owners close the loop")

    loop = [
        (0.5, 2.0, "Engine posts\ncards as PENDING", TEAL_LT, TEAL),
        (3.7, 2.0, "Owner sees brief\nin admin UI", BLUE_LT, BLUE),
        (6.9, 2.0, "Owner Accepts /\nDismisses / Done", ORANGE_LT, ORANGE),
        (10.1, 2.0, "Engine reads status\nbefore next run", PURPLE_LT, PURPLE),
    ]
    for left, top, text, fill, line in loop:
        sh = add_rounded_rect(s, Inches(left), top, Inches(2.7), Inches(2.0), fill, line=line, line_w=2)
        label_shape(sh, text, size=14, bold=True, color=DARK)

    for x in (3.25, 6.45, 9.65):
        add_arrow(s, Inches(x), Inches(3.0), Inches(x + 0.4), Inches(3.0), BLUE)

    add_text_box(
        s,
        Inches(0.5),
        Inches(4.6),
        Inches(12.3),
        Inches(2.0),
        "Simple product truth\n"
        "The AI proposes. The owner decides. The decision is stored on the platform so the next overnight run\n"
        "does not blindly re-suggest things the owner already dismissed.",
        size=15,
        color=DARK,
    )

    # ========== 14 Billing (light) ==========
    s = blank_slide(prs)
    slide_header(s, "Billing at a glance (still part of the story)", "Commercial system of record stays on the platform")

    add_bullet_box(
        s,
        Inches(0.5),
        Inches(1.4),
        Inches(6.0),
        Inches(5.0),
        "How charging works",
        [
            "One Stripe subscription per gas tenant",
            "Quantity = locations marked included_in_subscription",
            "Graduated pricing on a single Stripe Price",
            "Checkout + Customer Portal from admin billing",
            "Manual sync button pushes quantity to Stripe",
            "Auto webhook sync is a later platform step",
        ],
        fill=GREEN_LT,
        title_color=GREEN,
    )
    add_bullet_box(
        s,
        Inches(6.8),
        Inches(1.4),
        Inches(6.0),
        Inches(5.0),
        "Why developers care",
        [
            "Billing lives next to stations — not a third product",
            "Engine never talks to Stripe",
            "Demo UI works with Stripe test mode",
            "Keep commercial data out of model training",
            "Location CRUD can change billable quantity",
            "Keeps SaaS boundaries clean",
        ],
        fill=BLUE_LT,
        title_color=NAVY,
    )

    # ========== 15 Local / phases ==========
    s = blank_slide(prs)
    slide_header(s, "How to build & ship", "UI first with demo data; engine next")

    add_bullet_box(
        s,
        Inches(0.4),
        Inches(1.35),
        Inches(6.1),
        Inches(5.2),
        "Local developer path",
        [
            "Run gas station migration + seed SQL",
            "Point tenant at demo_gas_station_001",
            "Flip enable_gas_station_module = true",
            "npm run dev → open /admin/gas-station",
            "No AI engine required for UI work",
            "When building the engine: write-back to local REST API with service JWT",
        ],
        fill=BLUE_LT,
        title_color=NAVY,
    )
    add_bullet_box(
        s,
        Inches(6.8),
        Inches(1.35),
        Inches(6.1),
        Inches(5.2),
        "Phases (engine)",
        [
            "0 — repo + JWT smoke test",
            "1 — Verifone + Gilbarco files + rules brief",
            "2 — ML + LLM + nightly orchestration",
            "3 — chat + on-demand refresh",
            "4 — more vendors + hardening + alerts",
            "Pilot quality: roughly 10–15 person-weeks",
        ],
        fill=TEAL_LT,
        title_color=TEAL,
    )

    # ========== 16 Summary ==========
    s = blank_slide(prs)
    slide_header(s, "Takeaways for developers", "If you remember five things…")

    takeaways = [
        "Platform owns experience + billing + curated truth.",
        "AI engine owns connectors, models, and language.",
        "They meet only through typed DTOs + service JWT write-back.",
        "Raw POS never sits in the platform database.",
        "Build and demo the UI with seed data before any engine is live.",
    ]
    for i, t in enumerate(takeaways):
        top = Inches(1.3 + i * 1.0)
        card = add_rounded_rect(s, Inches(0.6), top, Inches(12.1), Inches(0.85), WHITE, line=BLUE, line_w=1.5)
        label_shape(card, f"{i + 1}.  {t}", size=16, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

    # ========== 17 Close ==========
    s = blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    add_text_box(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(0.8), "Build the brief. Write it back. Let owners click Accept.", size=28, bold=True, color=WHITE)
    add_text_box(
        s,
        Inches(0.8),
        Inches(3.3),
        Inches(11.5),
        Inches(2.0),
        "Source of truth for this deck\n"
        "• videos/gas_station_video_scripts_developer.html\n"
        "• gas_station_site_feasibility.md\n"
        "• gas_station_ai_engine_prd.html  ·  gas_station_requirements_implementation.html\n"
        "• Shared types: src/types/gasStation.ts",
        size=14,
        color=BLUE_LT,
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}")
    return OUT


if __name__ == "__main__":
    build()
